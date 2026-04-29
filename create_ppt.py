import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # Define colors
    bg_color = RGBColor(250, 246, 240)      # Soft beige
    primary_text = RGBColor(50, 45, 45)     # Dark brown/grey
    secondary_text = RGBColor(90, 85, 80)   # Lighter grey
    accent_color = RGBColor(180, 160, 140)  # Muted brown/taupe
    
    # We will use the blank slide layout for custom styling
    blank_slide_layout = prs.slide_layouts[6]
    
    def add_custom_slide(title_text=""):
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Add a subtle top border/gradient effect (solid rectangle here)
        top_bar = slide.shapes.add_shape(
            1, 0, 0, prs.slide_width, Inches(0.2) # 1 is msoShapeRectangle
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = accent_color
        top_bar.line.color.rgb = accent_color
        
        if title_text:
            # Add Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(1))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title_text.upper()
            p.font.bold = True
            p.font.size = Pt(36)
            p.font.name = "Calibri Light"
            p.font.color.rgb = primary_text
            
        return slide

    def add_card(slide, left, top, width, height, title, body_lines):
        # Background card
        card = slide.shapes.add_shape(
            1, left, top, width, height # 1 is msoShapeRectangle
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255) # White card
        card.line.color.rgb = RGBColor(230, 225, 220)
        
        # Text box over card
        txBox = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), height - Inches(0.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.name = "Calibri"
        p.font.color.rgb = primary_text
        
        for line in body_lines:
            p2 = tf.add_paragraph()
            p2.text = "• " + line
            p2.font.size = Pt(16)
            p2.font.name = "Calibri Light"
            p2.font.color.rgb = secondary_text
            p2.space_before = Pt(6)

    def add_body_text(slide, left, top, width, height, text, font_size=24, is_bold=False):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.bold = is_bold
        p.font.size = Pt(font_size)
        p.font.name = "Calibri Light"
        p.font.color.rgb = primary_text
        return tf
        
    def add_bullet_point(tf, text, font_size=20):
        p = tf.add_paragraph()
        p.text = "• " + text
        p.font.size = Pt(font_size)
        p.font.name = "Calibri Light"
        p.font.color.rgb = secondary_text
        p.space_before = Pt(8)

    # ---------------------------------------------------------
    # 1. TITLE SLIDE
    # ---------------------------------------------------------
    slide1 = add_custom_slide()
    # Center title
    txBox1 = slide1.shapes.add_textbox(Inches(1), Inches(2.5), prs.slide_width - Inches(2), Inches(2))
    tf1 = txBox1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "SMART ALUMNI MENTORING SYSTEM"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.bold = True
    p1.font.size = Pt(44)
    p1.font.name = "Calibri"
    p1.font.color.rgb = primary_text
    
    p2 = tf1.add_paragraph()
    p2.text = "Bridging the Gap Between Students and Industry Experts"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(24)
    p2.font.name = "Calibri Light"
    p2.font.color.rgb = accent_color
    p2.space_before = Pt(14)
    
    p3 = tf1.add_paragraph()
    p3.text = "\nPresented by: [Your Name]"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(18)
    p3.font.name = "Calibri Light"
    p3.font.color.rgb = secondary_text

    # ---------------------------------------------------------
    # 2. WHY THIS SYSTEM MATTERS
    # ---------------------------------------------------------
    slide2 = add_custom_slide("Why This System Matters")
    
    # Left Card: Problem
    add_card(slide2, Inches(1), Inches(2), Inches(3.8), Inches(4), "THE PROBLEM", [
        "Lack of direct mentorship",
        "No structured alumni connection",
        "Limited access to job referrals"
    ])
    
    # Right Card: Solution
    add_card(slide2, Inches(5.2), Inches(2), Inches(3.8), Inches(4), "THE SOLUTION", [
        "Dedicated alumni-student platform",
        "Verified and experienced mentors",
        "Structured scheduling and interaction"
    ])

    # ---------------------------------------------------------
    # 3. KEY BENEFITS
    # ---------------------------------------------------------
    slide3 = add_custom_slide("Key Benefits")
    
    w = Inches(2.8)
    add_card(slide3, Inches(0.5), Inches(2.5), w, Inches(3.5), "Career Guidance", [
        "Resume reviews",
        "Interview prep",
        "Industry insights"
    ])
    
    add_card(slide3, Inches(3.6), Inches(2.5), w, Inches(3.5), "Networking", [
        "Build professional bonds",
        "Direct referrals",
        "Community support"
    ])
    
    add_card(slide3, Inches(6.7), Inches(2.5), w, Inches(3.5), "Time Efficiency", [
        "Automated scheduling",
        "Centralized messaging",
        "Clear availability"
    ])

    # ---------------------------------------------------------
    # 4. SYSTEM ARCHITECTURE
    # ---------------------------------------------------------
    slide4 = add_custom_slide("System Architecture")
    
    add_card(slide4, Inches(1), Inches(2.5), Inches(2.5), Inches(3), "Frontend", [
        "HTML5",
        "Vanilla CSS (Modern UI)",
        "JavaScript (Dynamic)"
    ])
    
    add_card(slide4, Inches(3.75), Inches(2.5), Inches(2.5), Inches(3), "Backend", [
        "Node.js",
        "Express Framework",
        "RESTful APIs"
    ])
    
    add_card(slide4, Inches(6.5), Inches(2.5), Inches(2.5), Inches(3), "Database", [
        "MongoDB (NoSQL)",
        "Mongoose ODM",
        "Cloud Atlas"
    ])

    # ---------------------------------------------------------
    # 5. SYSTEM OVERVIEW
    # ---------------------------------------------------------
    slide5 = add_custom_slide("System Overview")
    tf5 = add_body_text(slide5, Inches(1), Inches(2), Inches(8), Inches(4), "A comprehensive platform to facilitate growth:", 24, True)
    
    add_bullet_point(tf5, "Mentorship Matching: Connect based on skills and interests.", 20)
    add_bullet_point(tf5, "Meeting Scheduling: Seamless calendar integration for 1-on-1s.", 20)
    add_bullet_point(tf5, "Job Board: Exclusive opportunities posted directly by alumni.", 20)
    add_bullet_point(tf5, "Integrated Chat: Real-time communication environment.", 20)

    # ---------------------------------------------------------
    # 6. STUDENT DASHBOARD
    # ---------------------------------------------------------
    slide6 = add_custom_slide("Student Dashboard")
    
    add_card(slide6, Inches(1), Inches(2), Inches(3.8), Inches(2), "Profile Management", [
        "Track CGPA, skills, projects"
    ])
    add_card(slide6, Inches(5.2), Inches(2), Inches(3.8), Inches(2), "Mentor Discovery", [
        "Find and request mentors"
    ])
    add_card(slide6, Inches(1), Inches(4.5), Inches(3.8), Inches(2), "Meeting Hub", [
        "Manage upcoming sessions"
    ])
    add_card(slide6, Inches(5.2), Inches(4.5), Inches(3.8), Inches(2), "Opportunities & Support", [
        "Browse jobs & file complaints"
    ])

    # ---------------------------------------------------------
    # 7. ALUMNI DASHBOARD
    # ---------------------------------------------------------
    slide7 = add_custom_slide("Alumni Dashboard")
    
    add_card(slide7, Inches(1.5), Inches(2.5), Inches(3.2), Inches(3), "Mentorship", [
        "Accept or reject requests",
        "Manage student roster"
    ])
    
    add_card(slide7, Inches(5.3), Inches(2.5), Inches(3.2), Inches(3), "Contribution", [
        "Provide meeting time slots",
        "Post exclusive job openings"
    ])

    # ---------------------------------------------------------
    # 8. ADMIN DASHBOARD
    # ---------------------------------------------------------
    slide8 = add_custom_slide("Admin Dashboard")
    
    tf8 = add_body_text(slide8, Inches(1), Inches(2), Inches(8), Inches(4), "Centralized Control & Moderation:", 24)
    add_bullet_point(tf8, "User Management: Oversee student and alumni accounts.")
    add_bullet_point(tf8, "Alumni Verification: Approve alumni registrations for security.")
    add_bullet_point(tf8, "Meeting Monitoring: Track system activity and engagement.")
    add_bullet_point(tf8, "Complaint Resolution: Address and resolve user issues.")

    # ---------------------------------------------------------
    # 9. MEETING SYSTEM (CORE FEATURE)
    # ---------------------------------------------------------
    slide9 = add_custom_slide("Meeting System Flow")
    
    # We create a simple visual flow
    steps = ["Student", "Request", "Alumni", "Provides Slots", "Student Selects", "Scheduled"]
    
    for i, step in enumerate(steps):
        left_pos = Inches(0.5 + (i * 1.5))
        top_pos = Inches(3.5) if i % 2 == 0 else Inches(2.5)
        
        card = slide9.shapes.add_shape(1, left_pos, top_pos, Inches(1.4), Inches(1))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color if i % 2 != 0 else RGBColor(255,255,255)
        card.line.color.rgb = accent_color
        
        txBox = slide9.shapes.add_textbox(left_pos, top_pos, Inches(1.4), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        p.text = step
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = primary_text

    # ---------------------------------------------------------
    # 10. CHAT SYSTEM
    # ---------------------------------------------------------
    slide10 = add_custom_slide("Chat System")
    
    add_card(slide10, Inches(1), Inches(2.5), Inches(2.5), Inches(3), "Student ↔ Alumni", [
        "Direct mentorship",
        "Secure messaging"
    ])
    add_card(slide10, Inches(3.75), Inches(2.5), Inches(2.5), Inches(3), "Student ↔ Admin", [
        "Support queries",
        "Issue resolution"
    ])
    add_card(slide10, Inches(6.5), Inches(2.5), Inches(2.5), Inches(3), "Alumni ↔ Admin", [
        "Platform feedback",
        "Verification help"
    ])

    # ---------------------------------------------------------
    # 11. KEY FEATURES
    # ---------------------------------------------------------
    slide11 = add_custom_slide("Key Features At A Glance")
    
    features = [
        "Role-Based Access", "Mentor Discovery", 
        "Meeting Scheduling", "Real-Time Chat", 
        "Job Sharing Board", "Complaint System"
    ]
    
    for i, feature in enumerate(features):
        row = i // 3
        col = i % 3
        
        left = Inches(1 + (col * 2.8))
        top = Inches(2.5 + (row * 1.8))
        
        # Simple card
        card = slide11.shapes.add_shape(1, left, top, Inches(2.5), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(230, 225, 220)
        
        txBox = slide11.shapes.add_textbox(left, top, Inches(2.5), Inches(1.2))
        p = txBox.text_frame.paragraphs[0]
        p.text = feature
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = primary_text

    # ---------------------------------------------------------
    # 12. CHALLENGES
    # ---------------------------------------------------------
    slide12 = add_custom_slide("Development Challenges")
    
    add_card(slide12, Inches(1), Inches(2.5), Inches(3.8), Inches(3), "Technical", [
        "Data handling & MongoDB queries",
        "State synchronization & real-time updates"
    ])
    
    add_card(slide12, Inches(5.2), Inches(2.5), Inches(3.8), Inches(3), "Design & Security", [
        "Maintaining UI consistency",
        "Securing routes and authentication"
    ])

    # ---------------------------------------------------------
    # 13. FUTURE SCOPE
    # ---------------------------------------------------------
    slide13 = add_custom_slide("Future Scope")
    
    tf13 = add_body_text(slide13, Inches(1), Inches(2.5), Inches(8), Inches(4), "Expanding the horizon:", 24)
    add_bullet_point(tf13, "AI Recommendations: Smart mentor-student matching algorithms.")
    add_bullet_point(tf13, "Video Integration: In-app meetings via Zoom or WebRTC.")
    add_bullet_point(tf13, "Mobile Application: Cross-platform access for on-the-go connectivity.")
    add_bullet_point(tf13, "Advanced Analytics: Tracking engagement and success metrics.")

    # ---------------------------------------------------------
    # 14. CONCLUSION
    # ---------------------------------------------------------
    slide14 = add_custom_slide("Conclusion")
    
    tf14 = add_body_text(slide14, Inches(1), Inches(2.5), Inches(8), Inches(4), "Empowering the Next Generation", 28, True)
    add_bullet_point(tf14, "Bridges the gap between academic learning and industry realities.", 22)
    add_bullet_point(tf14, "Fosters a strong, self-sustaining community of continuous learning.", 22)
    add_bullet_point(tf14, "Transforms connections into actionable career opportunities.", 22)

    prs.save("Smart_Alumni_Mentoring_System.pptx")

if __name__ == '__main__':
    create_presentation()
